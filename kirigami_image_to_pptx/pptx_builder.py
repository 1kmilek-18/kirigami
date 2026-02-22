"""
PPTX 生成（Req: 5.1–5.6）

レイヤー画像とテキスト要素から編集可能な .pptx を生成する。
レイヤーは Z-order 順に配置し、テキストは最前面にテキストボックスで配置する。
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import TextElement


def _get_image_size(path: str | Path) -> tuple[int, int]:
    """画像の幅・高さをピクセルで返す。"""
    from PIL import Image

    with Image.open(path) as img:
        return img.size


def _bbox_to_inches(
    bbox: tuple[float, float, float, float],
    ref_width: float,
    ref_height: float,
    slide_width_inches: float,
    slide_height_inches: float,
) -> tuple[float, float, float, float]:
    """
    bbox をスライド上のインチ座標に変換する。
    bbox が 0–1 正規化の場合は ref_width/ref_height を 1 としてスライド寸法でスケール、
    ピクセルの場合は ref に画像サイズを渡してスケールする。
    """
    x_min, y_min, x_max, y_max = bbox
    if max(x_min, y_min, x_max, y_max) <= 1.01:
        # 0–1 正規化座標
        scale_x, scale_y = slide_width_inches, slide_height_inches
    else:
        scale_x = slide_width_inches / ref_width if ref_width > 0 else 0
        scale_y = slide_height_inches / ref_height if ref_height > 0 else 0
    left = x_min * scale_x
    top = y_min * scale_y
    width = max(0.1, (x_max - x_min) * scale_x)
    height = max(0.1, (y_max - y_min) * scale_y)
    return (left, top, width, height)


def build_pptx(
    layer_paths: list[str],
    text_elements: list[TextElement],
    output_path: str | Path,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
    image_width: int | None = None,
    image_height: int | None = None,
) -> str:
    """
    レイヤー画像とテキスト要素から .pptx を生成する。

    Args:
        layer_paths: レイヤー画像のファイルパスリスト（Z-order 順、先頭が背景）
        text_elements: テキスト要素のリスト（bbox はピクセルまたは 0–1 正規化）
        output_path: 出力 .pptx パス
        slide_width_inches: スライド幅（インチ）。16:9 は 13.333
        slide_height_inches: スライド高さ（インチ）。16:9 は 7.5
        image_width: 参照画像幅（bbox がピクセル時のスケール用）。None の場合は先頭レイヤーから取得
        image_height: 参照画像高さ。None の場合は先頭レイヤーから取得

    Returns:
        出力したファイルパス（output_path）
    """
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    blank = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(blank)

    ref_w = float(image_width or 1)
    ref_h = float(image_height or 1)
    if layer_paths:
        ref_w, ref_h = _get_image_size(layer_paths[0])

    # レイヤーを Z-order 順に add_picture（先頭が背面）
    for path in layer_paths:
        p = Path(path)
        if not p.exists():
            continue
        slide.shapes.add_picture(
            str(p),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # テキストを最前面に add_textbox
    for el in text_elements:
        bbox = el.get("bbox")
        if not bbox or len(bbox) != 4:
            continue
        left, top, width, height = _bbox_to_inches(
            bbox, ref_w, ref_h, slide_width_inches, slide_height_inches
        )
        txb = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.get("text") or ""
        font = p.font
        if el.get("font_size_pt"):
            font.size = Pt(el["font_size_pt"])
        if el.get("font_color"):
            try:
                from pptx.dml.color import RGBColor

                hex_c = el["font_color"].lstrip("#")
                if len(hex_c) >= 6:
                    font.color.rgb = RGBColor(
                        int(hex_c[0:2], 16),
                        int(hex_c[2:4], 16),
                        int(hex_c[4:6], 16),
                    )
            except Exception:
                pass
        if el.get("is_bold") is True:
            font.bold = True

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)


def _add_slide(
    prs: Presentation,
    layer_paths: list[str],
    text_elements: list[TextElement],
    slide_width_inches: float,
    slide_height_inches: float,
    image_width: int | None,
    image_height: int | None,
) -> None:
    """1 スライドを追加する（build_pptx のコア部分）。"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    ref_w = float(image_width or 1)
    ref_h = float(image_height or 1)
    if layer_paths:
        ref_w, ref_h = _get_image_size(layer_paths[0])
    for path in layer_paths:
        p = Path(path)
        if not p.exists():
            continue
        slide.shapes.add_picture(
            str(p), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height
        )
    for el in text_elements:
        bbox = el.get("bbox")
        if not bbox or len(bbox) != 4:
            continue
        left, top, width, height = _bbox_to_inches(
            bbox, ref_w, ref_h, slide_width_inches, slide_height_inches
        )
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.get("text") or ""
        font = p.font
        if el.get("font_size_pt"):
            font.size = Pt(el["font_size_pt"])
        if el.get("font_color"):
            try:
                from pptx.dml.color import RGBColor
                hex_c = el["font_color"].lstrip("#")
                if len(hex_c) >= 6:
                    font.color.rgb = RGBColor(int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16))
            except Exception:
                pass
        if el.get("is_bold") is True:
            font.bold = True


def build_pptx_multi_slides(
    pages: list[tuple[list[str], list[TextElement]]],
    output_path: str | Path,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
) -> str:
    """
    複数ページ分のレイヤー・テキストから 1 つの .pptx を生成する（1 スライド = 1 ページ）。
    PDF 入力で 1 ファイル 1 .pptx とする場合に使用する。
    """
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    for layer_paths, text_elements in pages:
        ref_w, ref_h = None, None
        if layer_paths:
            ref_w, ref_h = _get_image_size(layer_paths[0])
        _add_slide(
            prs, layer_paths, text_elements,
            slide_width_inches, slide_height_inches, ref_w, ref_h,
        )
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
